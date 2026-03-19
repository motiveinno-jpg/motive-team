#!/usr/bin/env python3
"""토스페이먼츠 결제경로 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

SCREENSHOT_DIR = os.path.expanduser("~/motive-team/toss-ppt-screenshots")
OUTPUT_PATH = os.path.expanduser("~/motive-team/whistle-ai-결제경로.pptx")

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BLUE = RGBColor(0x21, 0x96, 0xF3)
RED = RGBColor(0xFF, 0x00, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)


def add_red_box(slide, left, top, width, height, line_width=Pt(3)):
    """Add a red rectangle outline (annotation box)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.background()
    shape.line.color.rgb = RED
    shape.line.width = line_width
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_slide_title(slide, number, title):
    """Add slide number and title at the top."""
    # Title background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"  {number}  "
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE

    run2 = p.add_run()
    run2.text = title
    run2.font.size = Pt(22)
    run2.font.bold = True
    run2.font.color.rgb = WHITE


def set_slide_bg(slide, color=None):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color or RGBColor(0xF5, 0xF5, 0xF5)


def add_screenshot(slide, img_path, left, top, max_width, max_height):
    """Add a screenshot image, fitting within bounds while maintaining aspect ratio."""
    img = Image.open(img_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Calculate dimensions to fit
    fit_w = max_width
    fit_h = int(fit_w / aspect)
    if fit_h > max_height:
        fit_h = max_height
        fit_w = int(fit_h * aspect)

    pic = slide.shapes.add_picture(img_path, left, top, fit_w, fit_h)
    return pic, fit_w, fit_h


def create_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ================================================================
    # SLIDE 1: 가맹점 정보 기재
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)
    add_slide_title(slide, "①", "가맹점 정보 기재")

    info_data = [
        ("상호명", "(주)모티브이노베이션"),
        ("사업자등록번호", "155-88-02209"),
        ("서비스 URL", "https://whistle-ai.com"),
        ("서비스명", "Whistle AI (AI 수출 통합 관리 플랫폼)"),
        ("토스 MID", "gmotivfyvh"),
        ("테스트 ID", "toss-review@whistle-ai.com"),
        ("테스트 PW", "TossReview2026!"),
    ]

    # Info table
    table_top = Inches(1.3)
    table_left = Inches(1.5)
    table_w = Inches(10)
    row_h = Inches(0.55)

    for i, (label, value) in enumerate(info_data):
        y = table_top + Emu(int(row_h * i))

        # Label cell background
        label_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, table_left, y, Inches(2.5), row_h
        )
        label_shape.fill.solid()
        label_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)
        label_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        label_shape.line.width = Pt(0.5)

        # Label text
        tf = label_shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Value cell background
        val_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, table_left + Inches(2.5), y,
            Inches(7.5), row_h
        )
        val_shape.fill.solid()
        val_shape.fill.fore_color.rgb = WHITE
        val_shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        val_shape.line.width = Pt(0.5)

        tf2 = val_shape.text_frame
        tf2.word_wrap = True
        tf2.paragraphs[0].alignment = PP_ALIGN.LEFT
        run2 = tf2.paragraphs[0].add_run()
        run2.text = f"  {value}"
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK_TEXT
        if label in ("테스트 ID", "테스트 PW"):
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(0xD3, 0x2F, 0x2F)

    # Note at bottom
    add_text_box(
        slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1),
        "※ 위 테스트 계정으로 로그인하시면 Starter 플랜 구독 상태의 서비스를 확인하실 수 있습니다.\n"
        "※ 결제 방식: 신용/체크카드 자동 정기결제 (빌링키 방식) | 결제 주기: 1개월 / 6개월 / 12개월",
        font_size=11, color=RGBColor(0x66, 0x66, 0x66)
    )

    # ================================================================
    # SLIDE 2: 하단정보 캡처
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "②", "하단정보 캡처 — whistle-ai.com/ko 푸터")

    img_path = os.path.join(SCREENSHOT_DIR, "02-footer.png")
    pic, pw, ph = add_screenshot(
        slide, img_path,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )

    # Red box around footer area (bottom of the screenshot)
    # Footer is at the very bottom of the image
    footer_box_top = Inches(1.1) + ph - Inches(1.0)
    add_red_box(slide, Inches(0.5), footer_box_top, pw, Inches(0.9))

    # Annotation label
    add_text_box(
        slide, Inches(0.5), Inches(1.1) + ph + Inches(0.1), Inches(12), Inches(0.5),
        "▲ 붉은 박스: 상호명 / 대표자명 / 사업자등록번호 / 통신판매업신고번호 / 주소 / 전화번호 / 이메일",
        font_size=11, bold=True, color=RED
    )

    # ================================================================
    # SLIDE 3: 환불규정 캡처
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "③", "환불규정 캡처 — whistle-ai.com/refund")

    img_path = os.path.join(SCREENSHOT_DIR, "03-refund.png")
    pic, pw, ph = add_screenshot(
        slide, img_path,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )

    # Red box around the refund policy title area
    add_red_box(slide, Inches(1.5), Inches(1.5), Inches(6), Inches(2.5))

    add_text_box(
        slide, Inches(0.5), Inches(1.1) + ph + Inches(0.1), Inches(12), Inches(0.5),
        "▲ 붉은 박스: 환불/취소 정책 (7일 이내 청약철회, 이용일수 일할 계산, SaaS 무형상품 환불규정)",
        font_size=11, bold=True, color=RED
    )

    # ================================================================
    # SLIDE 4: 로그인/회원가입 캡처
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "④", "로그인/회원가입 캡처 — whistle-ai.com/app")

    img_path = os.path.join(SCREENSHOT_DIR, "04-login.png")
    pic, pw, ph = add_screenshot(
        slide, img_path,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )

    # Red box around login form
    # Login form is centered, roughly in the middle
    form_left = Inches(4.0)
    form_top = Inches(2.0)
    add_red_box(slide, form_left, form_top, Inches(5.2), Inches(4.5))

    add_text_box(
        slide, Inches(0.5), Inches(1.1) + ph + Inches(0.1), Inches(12), Inches(0.5),
        "▲ 붉은 박스: 이메일+비밀번호 로그인 / Google 소셜 로그인 / 카카오 소셜 로그인 / 회원가입 탭",
        font_size=11, bold=True, color=RED
    )

    # ================================================================
    # SLIDE 5: 상품선택/구매과정 캡처
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "⑤", "상품선택/구매과정 캡처 — 로그인 후 앱 내 구독 & 서비스")

    img_path = os.path.join(SCREENSHOT_DIR, "05-subscription.png")
    pic, pw, ph = add_screenshot(
        slide, img_path,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )

    # Red box around pricing cards area
    add_red_box(slide, Inches(0.8), Inches(2.8), Inches(10.5), Inches(3.5))

    add_text_box(
        slide, Inches(0.5), Inches(1.1) + ph + Inches(0.1), Inches(12), Inches(0.5),
        "▲ 붉은 박스: 앱 내 구독 & 서비스 — AI분석 ₩25,000(단건) / Free ₩0 / Starter ₩99,000 / Pro ₩299,000 + 결제 주기 선택",
        font_size=11, bold=True, color=RED
    )

    # ================================================================
    # SLIDE 6: 카드 결제경로 캡처
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_slide_title(slide, "⑥", "카드 결제경로 캡처 — 토스페이먼츠 카드 등록 (빌링키)")

    img_path = os.path.join(SCREENSHOT_DIR, "06-toss-payment.png")
    pic, pw, ph = add_screenshot(
        slide, img_path,
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8)
    )

    # Red box around the Toss card registration form
    add_red_box(slide, Inches(3.0), Inches(1.5), Inches(7.0), Inches(5.0))

    add_text_box(
        slide, Inches(0.5), Inches(1.1) + ph + Inches(0.1), Inches(12), Inches(0.5),
        "▲ 붉은 박스: 토스페이먼츠 카드 등록 (빌링키 방식) — 카드번호 / 유효기간 / 주민번호 앞 6자리 입력 → 정기결제 자동 빌링",
        font_size=11, bold=True, color=RED
    )

    # Save
    prs.save(OUTPUT_PATH)
    print(f"PPT saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_ppt()
